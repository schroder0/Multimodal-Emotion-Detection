import collections
import collections.abc

# Fix for compatibility with python-pptx in newer python versions
collections.Mapping = collections.abc.Mapping
collections.Sequence = collections.abc.Sequence

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Setup Slide Layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
section_slide_layout = prs.slide_layouts[2]

# Helper for dark blue theme (Optional text color changes)
NAVY_BLUE = RGBColor(10, 22, 40)

def add_slide(title, points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = title
    tf = body_shape.text_frame
    
    for i, point in enumerate(points):
        if i == 0:
            tf.text = point[0]
            if len(point) > 1:
                p = tf.paragraphs[0]
                p.level = point[1]
        else:
            p = tf.add_paragraph()
            p.text = point[0]
            if len(point) > 1:
                p.level = point[1]

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(title_slide_layout)
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "Multimodal Emotion Detection via Chain-of-Thought Reasoning"
subtitle1.text = "Extending Generative Emotional Reasoning to Text, Emojis, and Images\n\nAaryan Kuntal"

# Slide 2: Introduction
add_slide("Introduction", [
    ("Emotion detection is a fundamental challenge in affective computing.", 0),
    ("Historically dominated by unimodal approaches and rigid classifiers.", 1),
    ("Human emotional expression is inherently multimodal.", 0),
    ("Utilizes text, facial expressions, vocal intonations, emojis, and images.", 1),
    ("Our approach: Leverage the generative capabilities of Large Language Models (LLMs).", 0),
    ("Integrate textual data, semantic emoji interpretation, and visual context.", 1)
])

# Slide 3: The Problem Statement
add_slide("The Problem Statement", [
    ("Why Multimodal? Why Reasoning?", 0),
    ("Traditional emotion detection uses rigid classification (e.g., text -> happy).", 1),
    ("Human emotion is complex and often masked.", 1),
    ("Saying 'I'm fine 🙂' while sitting alone in a dark room.", 2),
    ("We must look at text, emojis, and visual cues simultaneously to resolve ambiguity.", 1),
    ("The Goal: Move from classifying emotions to reasoning about them step-by-step.", 0)
])

# Slide 4: Reference Paper & Related Work
add_slide("Reference Paper & Related Work", [
    ("Base Paper: Emotion Detection via Chain-of-Thought Reasoning (arXiv:2408.04906v1)", 0),
    ("First work to use a generative approach to jointly address emotion detection and reasoning.", 1),
    ("Proves that asking an LLM to generate background knowledge improves accuracy.", 1),
    ("Our Extension:", 0),
    ("We take this text-based concept and expand it to a fully multimodal architecture.", 1),
    ("Integrates Vision-Language Models (BLIP) with LLMs (LLaMA).", 1)
])

# Slide 5: Datasets & Context
add_slide("Datasets & Context", [
    ("Designed to handle structures similar to standard benchmark datasets:", 0),
    ("CMU-MOSEI: Large-scale dataset for sentiment and emotion.", 1),
    ("IEMOCAP: Multimodal database of dyadic conversations.", 1),
    ("Our Experimental Setup:", 0),
    ("Curated evaluation set designed to test conflict-resolution.", 1),
    ("Congruent: Text matches image (e.g., Promotion + Celebration).", 2),
    ("Incongruent (Masking): Text contradicts image (e.g., 'I am fine' + Crying).", 2)
])

# Slide 6: System Architecture Flow
add_slide("System Architecture Flow", [
    ("End-to-End Processing Pipeline", 0),
    ("1. Input Collection: Frontend (React) gathers Text + Emoji + Image.", 1),
    ("2. Feature Extraction: Backend (FastAPI) splits modalities.", 1),
    ("3. Image Captioning: BLIP model translates image to text context.", 1),
    ("4. Emoji Semantic Mapping: Translates pictograms to emotional definitions.", 1),
    ("5. Context Assembly: Combines streams into a structured JSON prompt.", 1),
    ("6. Chain-of-Thought Reasoning: LLaMA 3.3 70B processes the prompt.", 1)
])

# Slide 7: The Models Used
add_slide("The Tech Stack under the Hood", [
    ("Vision Model: Salesforce/blip-image-captioning-base", 0),
    ("Performs zero-shot image captioning to describe the visual scene.", 1),
    ("Reasoning Engine: LLaMA 3.3 70B (via Groq API)", 0),
    ("Selected for exceptional zero-shot reasoning capabilities.", 1),
    ("Why Groq?: LPU architecture allows massive models to run in < 2 seconds.", 1),
    ("Why low temperature (0.3)?: Ensures analytical, deterministic reasoning.", 1)
])

# Slide 8: The 5-Step Chain-of-Thought
add_slide("5-Step Chain-of-Thought Prompting", [
    ("We force the LLM to 'think aloud' through 5 strict steps:", 0),
    ("Step 1: Text Analysis (Linguistic content and tone)", 1),
    ("Step 2: Emoji Analysis (Semantic meaning of emojis)", 1),
    ("Step 3: Image Analysis (Emotional context of the BLIP caption)", 1),
    ("Step 4: Conflict Resolution (Compare modalities and detect masking)", 1),
    ("Step 5: Final Prediction & Confidence Score", 1),
    ("This explicit structure dramatically reduces hallucinations.", 0)
])

# Slide 9: User Interface & Experience
add_slide("User Interface & Experience", [
    ("Modern, Client-Server Architecture", 0),
    ("Clean, bold Navy blue interface built with React.", 1),
    ("Supports drag-and-drop image uploads and emoji inputs.", 1),
    ("Explainable AI Output:", 0),
    ("The system doesn't just give an answer; it shows its work.", 1),
    ("Collapsible sections allow users to read the LLM's exact reasoning.", 1)
])

# Slide 10: Experimental Results
add_slide("Experimental Results", [
    ("Performance on Complex Modalities:", 0),
    ("Successfully resolves cross-modal conflicts (e.g., Sarcasm, Masking).", 1),
    ("Outperforms unimodal classifiers which default to the dominant text signal.", 1),
    ("Latency & Efficiency:", 0),
    ("End-to-end inference in under 2.5 seconds.", 1),
    ("Heuristic Fallback:", 0),
    ("Maintains a reliable baseline accuracy without LLM access using keyword weighting.", 1)
])

# Slide 11: Conclusion
add_slide("Conclusion", [
    ("Summary:", 0),
    ("Extending generative CoT reasoning to multimodal inputs yields a robust, accurate system.", 1),
    ("Explicit conflict-resolution prompts overcome limitations of rigid classifiers.", 1),
    ("High explainability builds trust in AI predictions.", 1),
    ("The Generative Paradigm:", 0),
    ("Reasoning step-by-step is fundamentally superior to direct classification for complex human emotions.", 1)
])

# Slide 12: Future Scope
add_slide("Future Scope", [
    ("Next Steps for Development:", 0),
    ("1. Native Multimodal LLMs (MLLMs):", 1),
    ("Migrating from BLIP + LLaMA to native models like LLaVA or GPT-4V to process images directly without captioning.", 2),
    ("2. Knowledge Distillation:", 1),
    ("Fine-tuning smaller, local models (e.g., LLaMA 3 8B) on the high-quality CoT data generated by the 70B model.", 2),
    ("3. Edge Deployment:", 1),
    ("Running the distilled models entirely on local hardware for privacy-preserving emotion detection.", 2)
])

# Save presentation
prs.save('/Users/aaryankuntal/Downloads/multimodal_emotion_detection/Presentation.pptx')
print("Successfully created Presentation.pptx with 12 detailed slides.")
